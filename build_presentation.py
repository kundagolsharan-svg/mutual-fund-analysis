"""Generate the slide deck for the Bluestock mutual fund capstone."""

from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_FILE = ROOT / "Bluestock_MF_Presentation.pptx"
IMAGE_PATHS = {
    "eda1": ROOT / "outputs" / "plots" / "nav_trends_all_schemes.png",
    "eda2": ROOT / "outputs" / "plots" / "monthly_sip_trend.png",
    "perf1": ROOT / "outputs" / "charts" / "benchmark_comparison.png",
    "perf2": ROOT / "outputs" / "charts" / "daily_returns_histograms.png",
    "dash1": ROOT / "outputs" / "plots" / "sector_allocation_donut.png",
    "dash2": ROOT / "outputs" / "plots" / "aum_by_house_year.png",
}


def add_bullets(shape, bullet_lines):
    text_frame = shape.text_frame
    text_frame.clear()
    for idx, line in enumerate(bullet_lines):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Calibri"


def add_slide(prs, title, body_lines=None, image_path=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if body_lines:
        add_bullets(slide.shapes.placeholders[1], body_lines)
    if image_path and image_path.exists():
        left = Inches(5.25)
        top = Inches(1.5)
        slide.shapes.add_picture(str(image_path), left, top, width=Inches(4.1))
    return slide


def add_image_slide(prs, title, subtitle, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), width=Inches(8))
    if subtitle:
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        tf = tx_box.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(14)
    return slide


def build_presentation():
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Bluestock MF Capstone"
    slide.placeholders[1].text = "Executive summary, performance analysis, and dashboard highlights"

    # Problem & objective
    add_slide(prs, "Problem & Objective", [
        "Assess mutual fund returns, risk, and benchmark performance across 40 schemes.",
        "Build a repeatable analytics pipeline for ETL, EDA, and performance reporting.",
        "Deliver a professional dashboard and final report for decision support.",
    ])

    # Data sources
    add_slide(prs, "Data Sources", [
        "Fund master, NAV history, AUM, SIP inflows, category inflows, folio counts, holdings, transactions, and benchmarks.",
        "Cleaned data is stored under data/processed for downstream analytics.",
    ])

    # Architecture
    add_slide(prs, "Architecture", [
        "Modular ETL with ingestion, cleaning, validation, analytics, and reporting steps.",
        "Processed CSV outputs and static dashboard artifacts support reproducible delivery.",
    ])

    # EDA highlight 1
    add_image_slide(prs, "EDA Highlight 1", "NAV trending across schemes reveals market cycles and fund movement.", IMAGE_PATHS["eda1"])

    # EDA highlight 2
    add_image_slide(prs, "EDA Highlight 2", "Monthly SIP inflows show strong retail contribution to fund flows.", IMAGE_PATHS["eda2"])

    # Performance metrics 1
    add_image_slide(prs, "Performance Metrics", "Top funds compared with NIFTY benchmarks to validate strategy performance.", IMAGE_PATHS["perf1"])

    # Performance metrics 2
    add_image_slide(prs, "Risk Metrics", "Daily return distributions and volatility are used to quantify downside exposure.", IMAGE_PATHS["perf2"])

    # Dashboard screenshot 1
    add_image_slide(prs, "Dashboard Screenshot", "Sector allocation and AUM distribution provide quick portfolio insight.", IMAGE_PATHS["dash1"])

    # Dashboard screenshot 2
    add_image_slide(prs, "Dashboard Screenshot", "AUM by fund house highlights concentration and market share.", IMAGE_PATHS["dash2"])

    # Key findings
    add_slide(prs, "Key Findings", [
        "Strong SIP inflows and AUM growth indicate continued investor confidence.",
        "Top-ranked funds show favorable risk-adjusted returns versus benchmarks.",
        "Correlation analysis underlines the importance of diversification across themes and sectors.",
    ])

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Thank You"
    add_bullets(slide.shapes.placeholders[1], [
        "Contact: Bluestock Analytics",
        "Repository includes final report, slide deck, and dashboard artifacts.",
    ])

    prs.save(str(OUTPUT_FILE))


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation saved to: {OUTPUT_FILE}")
